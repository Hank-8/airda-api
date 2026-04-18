"""
通用 PPTX 教案生成器
從 generate_color_sensor_pptx.py / generate_laser_sensor_pptx.py 提取通用邏輯
使用「送茶機器人.pptx」範本，根據 content dict 替換每頁內容
"""

import copy
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

TEMPLATE_DIR = os.path.join(os.path.dirname(__file__), 'templates')
TEMPLATE_PATH = os.path.join(TEMPLATE_DIR, '送茶機器人.pptx')


# ============================================================
# Utility functions (from existing scripts)
# ============================================================

def replace_text_preserve_format(shape, new_text):
    """Replace all text in a shape while preserving the first run's formatting."""
    if not hasattr(shape, 'text_frame'):
        return
    tf = shape.text_frame

    ref_font = None
    for para in tf.paragraphs:
        for run in para.runs:
            ref_font = run.font
            break
        if ref_font:
            break

    lines = new_text.split('\n')

    while len(tf.paragraphs) > 1:
        p_elem = tf.paragraphs[-1]._p
        p_elem.getparent().remove(p_elem)

    first_para = tf.paragraphs[0]
    first_para.clear()
    run = first_para.add_run()
    run.text = lines[0]
    if ref_font:
        _copy_font(ref_font, run.font)

    from pptx.oxml.ns import qn
    for line in lines[1:]:
        new_para = copy.deepcopy(first_para._p)
        for r_elem in new_para.findall(qn('a:r')):
            new_para.remove(r_elem)
        r_elem = copy.deepcopy(first_para._p.findall(qn('a:r'))[0])
        r_elem.find(qn('a:t')).text = line
        new_para.append(r_elem)
        tf._txBody.append(new_para)


def _copy_font(src, dst):
    if src.name:
        dst.name = src.name
    if src.size:
        dst.size = src.size
    if src.bold is not None:
        dst.bold = src.bold
    if src.italic is not None:
        dst.italic = src.italic
    try:
        if src.color and src.color.rgb:
            dst.color.rgb = src.color.rgb
    except AttributeError:
        pass


def remove_all_pictures(slide):
    to_remove = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    for s in to_remove:
        s._element.getparent().remove(s._element)
    return len(to_remove)


def resize_shape(shape, left=None, top=None, width=None, height=None):
    if left is not None: shape.left = left
    if top is not None: shape.top = top
    if width is not None: shape.width = width
    if height is not None: shape.height = height


def add_code_textbox(slide, code_text):
    left, top = Emu(914400), Emu(1371600)
    width, height = Emu(10363200), Emu(5200000)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code_text.strip().split('\n')):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_after = Pt(0)
        para.space_before = Pt(0)
        para.line_spacing = Pt(14)
        run = para.add_run()
        run.text = line
        run.font.name = 'Courier New'
        run.font.size = Pt(11)
        run.font.bold = False
        run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)


# ============================================================
# Generic slide processing
# ============================================================

def process_slides(prs, content: dict):
    """
    Process all 18 slides using the same keyword matching as the original scripts.
    content dict keys:
      cover_title, toc, objectives_knowledge, objectives_practice,
      theory_1_title, theory_1_body, theory_2_title, theory_2_body,
      theory_3_title, theory_3_body, theory_3_sub1, theory_3_sub2, theory_3_sub3,
      theory_4_title, theory_4_body,
      transition, hands_on, arduino_code,
      creative_title, creative_item1, creative_item2, creative_item3,
      teacher_notes, extension
    """
    slides = list(prs.slides)

    # --- Slide 1: Cover ---
    s = slides[0]
    for shape in s.shapes:
        if hasattr(shape, 'text_frame') and '送茶' in shape.text_frame.text:
            replace_text_preserve_format(shape, content['cover_title'])

    # --- Slide 2: Table of Contents ---
    s = slides[1]
    toc = content['toc']
    for shape in s.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        txt = shape.text_frame.text.strip()
        for old_key, new_val in toc.items():
            if txt.startswith(old_key):
                replace_text_preserve_format(shape, new_val)
                break

    # --- Slide 3: Classroom Rules --- (keep as-is)

    # --- Slide 4: Learning Objectives ---
    s = slides[3]
    remove_all_pictures(s)
    for shape in s.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        txt = shape.text_frame.text
        if '文化學習' in txt:
            resize_shape(shape, width=Emu(5029200), height=Emu(1600000))
            replace_text_preserve_format(shape, content['objectives_knowledge'])
        elif '創意組裝' in txt:
            resize_shape(shape, width=Emu(5029200), top=Emu(3962400))
            replace_text_preserve_format(shape, content['objectives_practice'])

    # --- Slide 5: Theory 1 ---
    s = slides[4]
    remove_all_pictures(s)
    for shape in s.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        txt = shape.text_frame.text
        if '服務生' in txt and len(txt) < 30:
            replace_text_preserve_format(shape, content['theory_1_title'])
        elif '侍應生' in txt or ('服務生' in txt and len(txt) > 30):
            replace_text_preserve_format(shape, content['theory_1_body'])

    # --- Slide 6: Theory 2 ---
    s = slides[5]
    remove_all_pictures(s)
    for shape in s.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        txt = shape.text_frame.text
        if '其他類型' in txt:
            replace_text_preserve_format(shape, content['theory_2_title'])
        elif '看板娘' in txt or '女服務生' in txt:
            resize_shape(shape, height=Emu(3600000))
            replace_text_preserve_format(shape, content['theory_2_body'])

    # --- Slide 7: Theory 3 ---
    s = slides[6]
    remove_all_pictures(s)
    for shape in s.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        txt = shape.text_frame.text
        if '主題服務生' in txt and len(txt) < 20:
            replace_text_preserve_format(shape, content['theory_3_title'])
        elif ('日本' in txt or '主題咖啡' in txt) and len(txt) > 20:
            resize_shape(shape, height=Emu(1600000))
            replace_text_preserve_format(shape, content['theory_3_body'])
        elif '不忍CAFE' in txt:
            replace_text_preserve_format(shape, content.get('theory_3_sub1', ''))
        elif '鐵道居酒屋' in txt:
            replace_text_preserve_format(shape, content.get('theory_3_sub2', ''))
        elif '謝幕' in txt:
            replace_text_preserve_format(shape, content.get('theory_3_sub3', ''))

    # --- Slide 8: Theory 4 ---
    s = slides[7]
    remove_all_pictures(s)
    for shape in s.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        txt = shape.text_frame.text
        if '小費' in txt and len(txt) < 20:
            replace_text_preserve_format(shape, content['theory_4_title'])
        elif '小費' in txt and len(txt) > 20:
            replace_text_preserve_format(shape, content['theory_4_body'])
        elif '相撲介紹' in txt or 'youtube' in txt.lower():
            resize_shape(shape, top=Emu(5715000), height=Emu(914400))
            replace_text_preserve_format(shape, content.get('theory_4_link', ''))

    # --- Slide 9: Transition ---
    s = slides[8]
    remove_all_pictures(s)
    for shape in s.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        txt = shape.text_frame.text
        if '送茶機器人' in txt or '小朋友' in txt:
            replace_text_preserve_format(shape, content['transition'])

    # --- Slide 10: Hands-on ---
    s = slides[9]
    remove_all_pictures(s)
    for shape in s.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        txt = shape.text_frame.text
        if 'reurl' in txt:
            resize_shape(shape,
                left=Emu(685800), top=Emu(1600000),
                width=Emu(5486400), height=Emu(4800000))
            replace_text_preserve_format(shape, content['hands_on'])

    # --- Slide 11: Timer --- (keep as-is)

    # --- Slide 12: Code Example ---
    s = slides[11]
    remove_all_pictures(s)
    add_code_textbox(s, content['arduino_code'])

    # --- Slide 13: Creative Challenge ---
    s = slides[12]
    for shape in s.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        txt = shape.text_frame.text
        if '創意改造' in txt and len(txt) < 10:
            replace_text_preserve_format(shape, content['creative_title'])
        elif '裝飾自己' in txt:
            replace_text_preserve_format(shape, content['creative_item1'])
        elif '修改程式' in txt:
            replace_text_preserve_format(shape, content['creative_item2'])
        elif '隨堂小任務' in txt:
            replace_text_preserve_format(shape, content['creative_item3'])

    # --- Slide 14-16: Clean Up / Share / Goodbye --- (keep as-is)

    # --- Slide 17: Teacher Notes ---
    s = slides[16]
    for shape in s.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        txt = shape.text_frame.text
        if '鼓勵學生' in txt:
            replace_text_preserve_format(shape, content['teacher_notes'])

    # --- Slide 18: Extended Discussion ---
    s = slides[17]
    for shape in s.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        txt = shape.text_frame.text
        if '組裝的作品' in txt or '招財貓' in txt:
            resize_shape(shape, height=Emu(3800000))
            replace_text_preserve_format(shape, content['extension'])


def generate_pptx(content: dict, output_path: str) -> str:
    """Load template, process slides, save to output_path."""
    prs = Presentation(TEMPLATE_PATH)
    process_slides(prs, content)
    prs.save(output_path)
    return output_path
